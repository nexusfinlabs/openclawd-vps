#!/usr/bin/env python3
"""
PPT Generator — Create professional presentations from natural language.

Usage:
  python3 ppt_generator.py "5 slides sobre AI en payments"
  python3 ppt_generator.py --template dark-gradient.pptx "5 slides sobre AI"

Templates:
  Place .pptx template files in ops/templates/ppt/
  The generator auto-detects slide layouts from the template.
  If no template is specified, uses built-in dark theme.

Flow:
  1. LLM generates structured slide content (JSON)
  2. python-pptx builds .pptx using template layouts or built-in styling
  3. Returns path to generated file
"""

import json
import os
import re
import sys
import glob
import requests
import logging
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Config ──────────────────────────────────────────────
DRAFTS_DIR = "/home/albi_agent/.openclaw/workspace/docs/drafts"
TEMPLATES_DIR = "/home/albi_agent/openclawd_stack/ops/templates/ppt"
OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY", "")

# ── Color Palette (built-in fallback) ──────────────────
COLORS = {
    "bg_dark": RGBColor(0x0B, 0x0F, 0x1A),
    "bg_card": RGBColor(0x13, 0x1A, 0x2E),
    "primary": RGBColor(0x38, 0xBD, 0xF8),
    "secondary": RGBColor(0x81, 0x8C, 0xF8),
    "accent": RGBColor(0x4A, 0xDE, 0x80),
    "warm": RGBColor(0xFB, 0xBF, 0x24),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "light_gray": RGBColor(0xCB, 0xD5, 0xE1),
    "text_muted": RGBColor(0x64, 0x74, 0x8B),
    "divider": RGBColor(0x1E, 0x29, 0x3B),
}

# ── Logging ─────────────────────────────────────────────
logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
log = logging.getLogger("ppt_generator")


def load_env():
    for path in ["/home/albi_agent/openclawd_stack/.env", ".env"]:
        if os.path.exists(path):
            with open(path) as f:
                for line in f:
                    if "=" in line and not line.startswith("#"):
                        k, v = line.strip().split("=", 1)
                        os.environ.setdefault(k, v.strip("\"'"))
    global OPENROUTER_API_KEY
    OPENROUTER_API_KEY = os.environ.get("OPENROUTER_API_KEY", "")


# ═══════════════════════════════════════════════════════
#  LLM Content Generation
# ═══════════════════════════════════════════════════════

def generate_slide_content(user_prompt):
    """Use LLM to generate structured slide content."""
    system = """Eres un diseñador de presentaciones de nivel ejecutivo.
Genera contenido para slides en JSON. Sé conciso y visual — máx 4-5 bullets por slide, texto corto.
Cada slide tiene: title, subtitle (opcional), bullets (lista, max 5), 
icon_hint (emoji sugerido para el tema), footer_note (opcional).

Tipos de slide:
- "cover": Portada. Solo title + subtitle.
- "content": Contenido principal con bullets.
- "two_column": Dos columnas (left_bullets + right_bullets, max 3 cada una).
- "highlight": Una frase grande destacada (quote field).
- "closing": Cierre con CTA. Solo title + subtitle.

Responde SOLO con un JSON array válido. Ejemplo:
[
  {"type": "cover", "title": "AI en Payments", "subtitle": "Tendencias & Oportunidades 2026", "icon_hint": "🚀"},
  {"type": "content", "title": "El Mercado Hoy", "bullets": ["$2.8T volumen global", "3x crecimiento LatAm", "Regulación acelerando"], "icon_hint": "📊"},
  {"type": "two_column", "title": "Riesgos vs Oportunidades", "left_title": "Riesgos", "left_bullets": ["Fragmentación", "Compliance"], "right_title": "Oportunidades", "right_bullets": ["Open Banking", "Instant payments"], "icon_hint": "⚖️"},
  {"type": "highlight", "quote": "El 68% de los pagos en LatAm serán digitales en 2028", "footer_note": "Fuente: McKinsey"},
  {"type": "closing", "title": "Next Steps", "subtitle": "dealflow@nexusfinlabs.com"}
]"""

    try:
        resp = requests.post(
            "https://openrouter.ai/api/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                "Content-Type": "application/json",
            },
            json={
                "model": "openai/gpt-4o-mini",
                "messages": [
                    {"role": "system", "content": system},
                    {"role": "user", "content": user_prompt},
                ],
                "temperature": 0.7,
            },
            timeout=30,
        )
        content = resp.json()["choices"][0]["message"]["content"]
        json_match = re.search(r'\[[\s\S]*\]', content)
        if json_match:
            return json.loads(json_match.group())
    except Exception as e:
        log.error("LLM slide generation failed: %s", e)
    return None


# ═══════════════════════════════════════════════════════
#  Template Detection
# ═══════════════════════════════════════════════════════

def list_available_templates():
    """List .pptx templates in the templates directory."""
    os.makedirs(TEMPLATES_DIR, exist_ok=True)
    return glob.glob(os.path.join(TEMPLATES_DIR, "*.pptx"))


def find_template(name=None):
    """Find a template by name or number. Returns None if no name given."""
    if not name:
        return None  # Don't auto-select — only use template when explicitly requested
    templates = list_available_templates()
    if not templates:
        return None
    # Try exact match first: "3" → "3.pptx"
    for t in templates:
        basename = os.path.basename(t)
        stem = os.path.splitext(basename)[0]
        if stem == name or basename.lower() == name.lower():
            return t
    # Fallback: substring match
    for t in templates:
        if name.lower() in os.path.basename(t).lower():
            return t
    return None


def detect_template_layouts(prs):
    """Detect and map slide layouts from a template."""
    layouts = {}
    for i, layout in enumerate(prs.slide_layouts):
        name = layout.name.lower()
        log.info("  Layout %d: '%s' (%d placeholders)", i, layout.name,
                 len(layout.placeholders))
        # Map common layout names
        if any(x in name for x in ["title slide", "portada", "cover", "título"]):
            layouts["cover"] = layout
        elif any(x in name for x in ["two content", "dos", "comparison", "two col"]):
            layouts["two_column"] = layout
        elif any(x in name for x in ["section", "sección", "divider"]):
            layouts["highlight"] = layout
        elif any(x in name for x in ["title and content", "título y obj", "content"]):
            layouts["content"] = layout
        elif any(x in name for x in ["blank", "en blanco"]):
            layouts["blank"] = layout
    return layouts


# ═══════════════════════════════════════════════════════
#  Built-in Slide Builders (no template)
# ═══════════════════════════════════════════════════════

def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _bar(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def _text(slide, text, x, y, w, h, size=18, color=None, bold=False,
          align=PP_ALIGN.LEFT, font="Calibri"):
    if color is None:
        color = COLORS["white"]
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def _add_bullets(slide, bullets, x_start, y_start, width, colors_cycle=None):
    """Add styled bullet points."""
    if colors_cycle is None:
        colors_cycle = [COLORS["primary"], COLORS["secondary"], COLORS["accent"], COLORS["warm"]]
    for i, bullet in enumerate(bullets[:6]):
        y = y_start + i * Inches(0.82)
        # Colored bar indicator
        _bar(slide, x_start, y + Inches(0.08), Inches(0.05), Inches(0.5),
             colors_cycle[i % len(colors_cycle)])
        # Text
        _text(slide, bullet, x_start + Inches(0.25), y, width, Inches(0.7),
              size=15, color=COLORS["light_gray"])


def build_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, COLORS["bg_dark"])

    # Decorative shapes
    _bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.05), COLORS["primary"])
    # Large circle decoration
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(-0.5), Inches(3.5), Inches(3.5))
    c.fill.solid()
    c.fill.fore_color.rgb = COLORS["bg_card"]
    c.line.fill.background()
    # Smaller accent circle
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.3), Inches(1.2), Inches(1.2))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLORS["primary"]
    c2.line.fill.background()

    # Icon hint
    icon = data.get("icon_hint", "📊")
    _text(slide, icon, Inches(0.8), Inches(1.8), Inches(1), Inches(0.8),
          size=40, align=PP_ALIGN.LEFT)

    # Accent line
    _bar(slide, Inches(0.8), Inches(2.8), Inches(1.2), Inches(0.05), COLORS["primary"])

    # Title
    _text(slide, data.get("title", "Presentación"), Inches(0.8), Inches(3.0),
          Inches(7.5), Inches(1.5), size=42, bold=True)

    # Subtitle
    if data.get("subtitle"):
        _text(slide, data["subtitle"], Inches(0.8), Inches(4.6),
              Inches(7.5), Inches(0.8), size=20, color=COLORS["light_gray"])

    # Branding
    _text(slide, "Alberto Lebrón · Nexus FinLabs", Inches(0.8), Inches(6.5),
          Inches(4), Inches(0.4), size=11, color=COLORS["text_muted"])
    _text(slide, datetime.now().strftime("%B %Y"), Inches(7), Inches(6.5),
          Inches(2.5), Inches(0.4), size=11, color=COLORS["text_muted"],
          align=PP_ALIGN.RIGHT)


def build_content(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, COLORS["bg_dark"])

    # Top bar
    _bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.04), COLORS["primary"])

    # Slide number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.5), Inches(0.35), Inches(0.55), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS["primary"]
    badge.line.fill.background()
    tf = badge.text_frame
    tf.text = f"{num:02d}"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = COLORS["bg_dark"]
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Icon + Title on same line
    icon = data.get("icon_hint", "")
    title_text = f"{icon}  {data.get('title', '')}" if icon else data.get("title", "")
    _text(slide, title_text, Inches(1.3), Inches(0.3), Inches(8), Inches(0.6),
          size=26, bold=True)

    # Accent line
    _bar(slide, Inches(1.3), Inches(1.0), Inches(1.2), Inches(0.04), COLORS["primary"])

    # Bullets with colored indicators
    _add_bullets(slide, data.get("bullets", []), Inches(1.3), Inches(1.3), Inches(7.5))

    # Footer note
    if data.get("footer_note"):
        _text(slide, data["footer_note"], Inches(0.8), Inches(6.6),
              Inches(8.4), Inches(0.3), size=9, color=COLORS["text_muted"])


def build_two_column(prs, data, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, COLORS["bg_dark"])

    _bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.04), COLORS["primary"])

    # Title
    icon = data.get("icon_hint", "")
    title_text = f"{icon}  {data.get('title', '')}" if icon else data.get("title", "")
    _text(slide, title_text, Inches(0.8), Inches(0.3), Inches(8.4), Inches(0.6),
          size=26, bold=True)
    _bar(slide, Inches(0.8), Inches(1.0), Inches(1.2), Inches(0.04), COLORS["primary"])

    # Left column card
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.6), Inches(1.3), Inches(4.1), Inches(5.2))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLORS["bg_card"]
    left_card.line.fill.background()

    left_title = data.get("left_title", "")
    if left_title:
        _text(slide, left_title, Inches(0.9), Inches(1.5), Inches(3.5), Inches(0.5),
              size=18, bold=True, color=COLORS["primary"])

    _add_bullets(slide, data.get("left_bullets", []), Inches(0.9), Inches(2.2),
                 Inches(3.5), [COLORS["primary"], COLORS["accent"]])

    # Right column card
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(5.3), Inches(1.3), Inches(4.1), Inches(5.2))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLORS["bg_card"]
    right_card.line.fill.background()

    right_title = data.get("right_title", "")
    if right_title:
        _text(slide, right_title, Inches(5.6), Inches(1.5), Inches(3.5), Inches(0.5),
              size=18, bold=True, color=COLORS["secondary"])

    _add_bullets(slide, data.get("right_bullets", []), Inches(5.6), Inches(2.2),
                 Inches(3.5), [COLORS["secondary"], COLORS["warm"]])


def build_highlight(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, COLORS["bg_dark"])

    _bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.05), COLORS["secondary"])

    # Central card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1), Inches(1.5), Inches(8), Inches(4.5))
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["bg_card"]
    card.line.fill.background()

    # Quote marks
    _text(slide, "❝", Inches(1.5), Inches(1.8), Inches(1), Inches(1),
          size=48, color=COLORS["primary"])

    # Quote text
    _text(slide, data.get("quote", data.get("title", "")),
          Inches(1.8), Inches(2.8), Inches(6.4), Inches(2),
          size=28, color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)

    if data.get("footer_note"):
        _text(slide, data["footer_note"], Inches(1), Inches(6.3),
              Inches(8), Inches(0.4), size=11, color=COLORS["text_muted"],
              align=PP_ALIGN.CENTER)

    _bar(slide, Inches(0), Inches(7.2), Inches(10), Inches(0.05), COLORS["primary"])


def build_closing(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, COLORS["bg_dark"])

    _bar(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), COLORS["secondary"])

    # Decorative circles
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(0.3), Inches(3), Inches(3))
    c.fill.solid()
    c.fill.fore_color.rgb = COLORS["bg_card"]
    c.line.fill.background()

    _text(slide, data.get("title", "Gracias"), Inches(1), Inches(2.5),
          Inches(8), Inches(1.2), size=44, bold=True, align=PP_ALIGN.CENTER)

    if data.get("subtitle"):
        _text(slide, data["subtitle"], Inches(1), Inches(4.0),
              Inches(8), Inches(0.8), size=18, color=COLORS["primary"],
              align=PP_ALIGN.CENTER)

    _text(slide, "dealflow@nexusfinlabs.com · +34 605 693 177",
          Inches(1), Inches(6.3), Inches(8), Inches(0.4),
          size=12, color=COLORS["text_muted"], align=PP_ALIGN.CENTER)

    _bar(slide, Inches(0), Inches(7.2), Inches(10), Inches(0.06), COLORS["primary"])


# ═══════════════════════════════════════════════════════
#  Template-Based Slide Builders
# ═══════════════════════════════════════════════════════

def _fill_placeholders(slide, data, slide_type="content"):
    """Fill placeholders in a template slide with generated content."""
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        # idx 0 = title, idx 1 = subtitle/body, idx 2+ = other
        if idx == 0:
            ph.text = data.get("title", "")
        elif idx == 1:
            if slide_type == "cover":
                ph.text = data.get("subtitle", "")
            elif slide_type == "highlight":
                ph.text = data.get("quote", data.get("title", ""))
            else:
                # Format bullets into the body placeholder
                bullets = data.get("bullets", [])
                if bullets:
                    ph.text = bullets[0]
                    tf = ph.text_frame
                    for bullet in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0
                else:
                    ph.text = ""
        elif idx == 2:
            # Second body / notes
            ph.text = data.get("footer_note", "")


def build_template_slide(prs, layouts, data, slide_type, slide_num):
    """Build a slide using template layouts."""
    layout = layouts.get(slide_type) or layouts.get("content") or layouts.get("blank")
    if not layout:
        # Fallback to first non-blank layout
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)
    _fill_placeholders(slide, data, slide_type)
    return slide


# ═══════════════════════════════════════════════════════
#  Main Generator
# ═══════════════════════════════════════════════════════

def generate_pptx(user_prompt, template_name=None):
    """Main function: LLM content → python-pptx → .pptx file."""
    load_env()
    log.info("Generating PPT for: %s", user_prompt[:80])

    # 1. Generate content via LLM
    slides_data = generate_slide_content(user_prompt)
    if not slides_data:
        return None, "❌ LLM no pudo generar el contenido de las slides"

    log.info("LLM generated %d slides", len(slides_data))

    # 2. Check for template
    template_path = find_template(template_name)
    using_template = False

    if template_path:
        log.info("Using template: %s", template_path)
        prs = Presentation(template_path)
        layouts = detect_template_layouts(prs)
        using_template = True
        # Remove any existing slides from the template
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        log.info("No template found — using built-in dark theme")
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        layouts = None

    # 3. Build slides
    slide_num = 0
    for sdata in slides_data:
        stype = sdata.get("type", "content")

        if using_template:
            build_template_slide(prs, layouts, sdata, stype, slide_num)
        else:
            if stype == "cover":
                build_cover(prs, sdata)
            elif stype == "closing":
                build_closing(prs, sdata)
            elif stype == "two_column":
                slide_num += 1
                build_two_column(prs, sdata, slide_num)
            elif stype == "highlight":
                build_highlight(prs, sdata)
            else:
                slide_num += 1
                build_content(prs, sdata, slide_num)

    # 4. Save
    os.makedirs(DRAFTS_DIR, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = re.sub(r'[^a-zA-Z0-9]', '_', user_prompt[:30]).strip('_')
    filename = f"PPT_{safe_name}_{timestamp}.pptx"
    filepath = os.path.join(DRAFTS_DIR, filename)

    prs.save(filepath)
    log.info("PPTX saved: %s", filepath)

    template_info = f"📐 Template: `{os.path.basename(template_path)}`\n" if using_template else ""
    summary = (
        f"📊 *Presentación PPTX generada*\n\n"
        f"📁 Archivo: `{filename}`\n"
        f"📄 Slides: {len(slides_data)}\n"
        f"{template_info}"
        f"📎 Ruta: {filepath}\n"
    )

    return filepath, summary


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="PPT Generator")
    parser.add_argument("prompt", nargs="*", help="Descripción de la presentación")
    parser.add_argument("--template", "-t", default=None, help="Nombre del template .pptx")
    parser.add_argument("--prompt-file", default=None, help="Read prompt from file (for large context)")
    args = parser.parse_args()

    if args.prompt_file:
        with open(args.prompt_file, "r", encoding="utf-8") as f:
            prompt = f.read().strip()
    elif args.prompt:
        prompt = " ".join(args.prompt)
    else:
        print("Usage: python3 ppt_generator.py \"descripción\" or --prompt-file /path")
        sys.exit(1)

    filepath, message = generate_pptx(prompt, template_name=args.template)
    print(message)
